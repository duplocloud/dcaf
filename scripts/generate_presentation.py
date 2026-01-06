#!/usr/bin/env python3
"""
Generate PowerPoint presentation for DCAF Core.

Usage:
    pip install python-pptx
    python scripts/generate_presentation.py

Output:
    docs/presentations/dcaf-core-presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(30, 30, 46)       # Dark background
ACCENT = RGBColor(137, 180, 250)      # Blue accent
TEXT_WHITE = RGBColor(205, 214, 244)  # Off-white text
TEXT_DIM = RGBColor(147, 153, 178)    # Dimmed text
CODE_BG = RGBColor(45, 45, 65)        # Code block background
GREEN = RGBColor(166, 227, 161)       # Success green


def create_presentation():
    """Create the DCAF Core presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
        "DCAF Core",
        "Building AI Agents with Simplicity",
        "DuploCloud Agent Framework"
    )
    
    # Slide 2: What is DCAF Core?
    add_content_slide(prs,
        "What is DCAF Core?",
        "A Simple API for Complex Agent Workflows",
        [
            "🛠️  Call tools – Execute functions on behalf of users",
            "✅  Request approval – Pause for human authorization",
            "📡  Stream responses – Real-time token-by-token output",
            "🌐  Serve via REST – One-line HTTP server deployment",
        ],
        footer="Philosophy: Hide complexity, expose simplicity."
    )
    
    # Slide 3: The Core API in 10 Lines
    add_code_slide(prs,
        "The Core API in 10 Lines",
        '''from dcaf.core import Agent, serve
from dcaf.tools import tool

@tool(description="List Kubernetes pods")
def list_pods(namespace: str = "default") -> str:
    return kubectl(f"get pods -n {namespace}")

agent = Agent(tools=[list_pods], system="You are a K8s assistant.")

serve(agent)''',
        footer="That's it. Your agent is now running at http://localhost:8000"
    )
    
    # Slide 4: Tool Definitions
    add_code_slide(prs,
        "Tool Definitions",
        '''# Auto-generated schemas from type hints
@tool(description="Delete a Kubernetes pod")
def delete_pod(name: str, namespace: str = "default") -> str:
    return kubectl(f"delete pod {name} -n {namespace}")

# Adding approval requirements
@tool(requires_approval=True, description="Delete a Kubernetes pod")
def delete_pod(name: str, namespace: str = "default") -> str:
    return kubectl(f"delete pod {name} -n {namespace}")''',
        footer="Schemas inferred from function signatures automatically"
    )
    
    # Slide 5: The Agent Class
    add_code_slide(prs,
        "The Agent Class",
        '''from dcaf.core import Agent

agent = Agent(
    tools=[list_pods, delete_pod, restart_deployment],
    system="You are a Kubernetes assistant.",
)

# Running the agent programmatically
response = agent.run([
    {"role": "user", "content": "What pods are running?"}
])

print(response.text)''',
        footer="Minimal configuration, maximum capability"
    )
    
    # Slide 6: The serve() Function
    add_code_slide(prs,
        "The serve() Function",
        '''from dcaf.core import serve

# One-line REST API server
serve(agent)  # Running at http://0.0.0.0:8000

# With configuration options
serve(
    agent,
    port=8000,              # Port to listen on
    host="0.0.0.0",         # Host to bind to
    reload=True,            # Auto-reload for development
    log_level="info",       # Logging verbosity
)''',
        footer="One function to expose your agent via HTTP"
    )
    
    # Slide 7: Production Configuration
    add_content_slide(prs,
        "Production Configuration",
        "Built-in Support for Production Deployments",
        [
            "workers=4           → Multiple worker processes",
            "timeout_keep_alive=30  → Match load balancer timeout",
            "log_level='warning'    → Reduce noise in production",
        ],
        footer="Best Practice: workers = (2 × cpu_cores) + 1"
    )
    
    # Slide 8: REST API Endpoints
    add_table_slide(prs,
        "REST API Endpoints",
        ["Endpoint", "Method", "Description"],
        [
            ["/health", "GET", "Health check (always responds immediately)"],
            ["/api/chat", "POST", "Synchronous chat"],
            ["/api/chat-stream", "POST", "Streaming chat (NDJSON)"],
        ]
    )
    
    # Slide 9: Human-in-the-Loop Approval
    add_content_slide(prs,
        "Human-in-the-Loop Approval",
        "How It Works",
        [
            "1. User asks agent to perform an action",
            "2. Agent identifies a tool that requires approval",
            "3. Agent pauses and returns pending tool calls",
            "4. User reviews and approves (or rejects)",
            "5. Agent executes the approved tool",
        ],
        footer="Built-in safety for dangerous operations"
    )
    
    # Slide 10: Custom Agent Logic
    add_code_slide(prs,
        "Custom Agent Logic",
        '''def my_agent(messages: list, context: dict) -> AgentResult:
    tenant = context.get("tenant_name")
    
    # Classify intent first
    classifier = Agent(system="Classify as: query or action")
    intent = classifier.run(messages)
    
    if "action" in intent.text.lower():
        executor = Agent(tools=[list_pods, delete_pod])
        result = executor.run(messages)
        return AgentResult(text=result.text)
    
    return AgentResult(text=intent.text)

serve(my_agent)''',
        footer="Any structure you need: multiple LLM calls, branching, orchestration"
    )
    
    # Slide 11: Platform Context
    add_code_slide(prs,
        "Platform Context",
        '''def my_agent(messages: list, context: dict) -> AgentResult:
    # Context is extracted automatically
    tenant = context.get("tenant_name")
    namespace = context.get("k8s_namespace")
    user_id = context.get("user_id")
    
    agent = Agent(
        tools=[...],
        system=f"Assisting tenant {tenant} in namespace {namespace}"
    )
    ...''',
        footer="Available: tenant_name, k8s_namespace, user_id, duplo_token, and more"
    )
    
    # Slide 12: Streaming Responses
    add_content_slide(prs,
        "Streaming Responses",
        "Real-Time Token-by-Token Output",
        [
            'text_delta    → Incremental text from the LLM',
            'tool_calls    → Tools requiring approval',
            'executed_tool_calls → Results from executed tools',
            'done          → Stream completed',
            'error         → An error occurred',
        ],
        footer="NDJSON format for easy parsing"
    )
    
    # Slide 13: Adding Custom Endpoints
    add_code_slide(prs,
        "Adding Custom Endpoints",
        '''from fastapi import APIRouter

custom_router = APIRouter()

@custom_router.get("/api/custom/schema")
async def get_schema():
    return {"tools": ["list_pods", "delete_pod"]}

@custom_router.get("/api/custom/health")
async def detailed_health():
    return {"status": "healthy", "tools_loaded": 2}

# Include custom routes
serve(agent, additional_routers=[custom_router])''',
        footer="Extend your agent server with any FastAPI routes"
    )
    
    # Slide 14: Docker Deployment
    add_code_slide(prs,
        "Docker Deployment",
        '''# main.py
import os
from dcaf.core import Agent, serve

agent = Agent(tools=[...], system="You are a K8s assistant.")

if __name__ == "__main__":
    serve(
        agent,
        host="0.0.0.0",
        port=int(os.getenv("PORT", 8000)),
        workers=int(os.getenv("WORKERS", 4)),
        timeout_keep_alive=int(os.getenv("KEEP_ALIVE", 30)),
    )''',
        footer="docker run -e WORKERS=8 -e KEEP_ALIVE=60 my-agent:latest"
    )
    
    # Slide 15: Complete Kubernetes Agent Example
    add_code_slide(prs,
        "Complete Example: Kubernetes Agent",
        '''@tool(description="List Kubernetes pods")
def list_pods(namespace: str = "default") -> str:
    return kubectl(f"get pods -n {namespace}")

@tool(requires_approval=True, description="Delete a pod")
def delete_pod(name: str, namespace: str = "default") -> str:
    return kubectl(f"delete pod {name} -n {namespace}")

@tool(requires_approval=True, description="Restart a deployment")
def restart_deployment(name: str, namespace: str = "default") -> str:
    return kubectl(f"rollout restart deployment {name} -n {namespace}")

agent = Agent(tools=[list_pods, delete_pod, restart_deployment])
serve(agent, port=8000, workers=4)''',
        footer="Production-ready in ~15 lines of code"
    )
    
    # Slide 16: Architecture Overview
    add_architecture_slide(prs)
    
    # Slide 17: Key Concepts
    add_table_slide(prs,
        "Key Concepts",
        ["Concept", "Description"],
        [
            ["Agent", "Your LLM-powered assistant with tools"],
            ["Tool", "A function the agent can call"],
            ["Approval", "Human authorization for sensitive tools"],
            ["serve()", "One-line REST API server"],
            ["create_app()", "Programmatic FastAPI control"],
            ["Platform Context", "Runtime environment (tenant, namespace, etc.)"],
            ["AgentResult", "Return type for custom agent functions"],
        ]
    )
    
    # Slide 18: Getting Started
    add_content_slide(prs,
        "Getting Started",
        "Installation",
        [
            "pip install git+https://github.com/duplocloud/dcaf.git",
            "",
            "Documentation:",
            "  • docs/core/index.md – Core Overview",
            "  • docs/core/server.md – Server Guide",
            "  • docs/guides/custom-agents.md – Custom Agents",
        ]
    )
    
    # Slide 19: Q&A
    add_qa_slide(prs)
    
    return prs


def add_title_slide(prs, title, subtitle, footer):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    add_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_DIM
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, subtitle, bullets, footer=None):
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    add_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_WHITE
    
    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.5), Inches(4.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(12)
    
    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = TEXT_DIM


def add_code_slide(prs, title, code, footer=None):
    """Add a slide with code."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    add_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Code background
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.0)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = CODE_BG
    code_bg.line.fill.background()
    
    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.6))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    lines = code.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(2)
    
    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = GREEN


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    add_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * num_rows)
    ).table
    
    # Style table
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_BG
    
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = CODE_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_WHITE


def add_architecture_slide(prs):
    """Add the architecture diagram slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    add_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture Overview"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Your Code box
    add_arch_box(slide, 1.5, 1.5, 10.3, 1.2, "Your Code", 
                 "agent = Agent(tools=[...])  →  serve(agent)", ACCENT)
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(6), Inches(2.8), Inches(1), Inches(0.5))
    tf = arrow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "▼"
    p.font.size = Pt(32)
    p.font.color.rgb = TEXT_DIM
    p.alignment = PP_ALIGN.CENTER
    
    # DCAF Core box
    add_arch_box(slide, 1.5, 3.3, 10.3, 1.5, "DCAF Core",
                 "HTTP → Messages → Agent Logic → Approvals → Response", GREEN)
    
    # Arrow
    arrow_box2 = slide.shapes.add_textbox(Inches(6), Inches(4.9), Inches(1), Inches(0.5))
    tf = arrow_box2.text_frame
    p = tf.paragraphs[0]
    p.text = "▼"
    p.font.size = Pt(32)
    p.font.color.rgb = TEXT_DIM
    p.alignment = PP_ALIGN.CENTER
    
    # LLM box
    add_arch_box(slide, 1.5, 5.4, 10.3, 1.0, "LLM (AWS Bedrock)",
                 "Claude 3.5 Sonnet / Claude 4 / etc.", TEXT_WHITE)


def add_arch_box(slide, left, top, width, height, title, content, color):
    """Add an architecture box."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = color
    box.line.width = Pt(2)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.4), Inches(0.5)
    )
    tf = content_box.text_frame
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER


def add_qa_slide(prs):
    """Add Q&A slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    add_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Key takeaways
    takeaways = [
        "Simple API – Agent + serve() is all you need",
        "Tool calling – Decorator-based with auto-generated schemas",
        "Human-in-the-loop – Built-in approval for sensitive operations",
        "Production-ready – Workers, keep-alive, health checks included",
        "Flexible – Use classes or functions, your choice",
    ]
    
    bullet_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(3))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(takeaways):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓  {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)


def add_background(slide, color):
    """Add solid color background to slide."""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def main():
    """Generate the presentation."""
    print("Generating DCAF Core presentation...")
    
    prs = create_presentation()
    
    # Ensure output directory exists
    output_dir = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations")
    os.makedirs(output_dir, exist_ok=True)
    
    output_path = os.path.join(output_dir, "dcaf-core-presentation.pptx")
    prs.save(output_path)
    
    print(f"✓ Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

