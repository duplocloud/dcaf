#!/usr/bin/env python3
"""Generate DCAF Core PowerPoint presentation with white background and black text."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def rgb_color(r, g, b):
    """Helper to create RGB color."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


# Colors
BLACK = rgb_color(0, 0, 0)
GRAY = rgb_color(100, 100, 100)
LIGHT_GRAY = rgb_color(245, 245, 245)


def add_title_slide(title, subtitle=None):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(title, subtitle=None):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(title, content_lines, code=None):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    top = 1.4
    if content_lines:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.333), Inches(3))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.color.rgb = BLACK
            p.space_before = Pt(8)
            p.space_after = Pt(4)
        top += len(content_lines) * 0.4 + 0.5
    
    if code:
        code_top = top if content_lines else 1.4
        code_height = min(len(code.split('\n')) * 0.3 + 0.4, 5.5)
        
        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.5), Inches(code_top), 
            Inches(12.333), Inches(code_height)
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = LIGHT_GRAY
        code_bg.line.fill.background()
        
        code_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(code_top + 0.15), 
            Inches(12), Inches(code_height - 0.3)
        )
        tf = code_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = BLACK
    
    return slide


def add_table_slide(title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLACK
    
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.4),
        Inches(12.333), Inches(num_rows * 0.5)
    ).table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = BLACK
    
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = BLACK
    
    return slide


# ============================================================================
# SECTION 1: INTRODUCTION
# ============================================================================

add_title_slide("DCAF Core", "Building AI Agents with Simplicity")

add_content_slide(
    "What is DCAF Core?",
    [
        "DCAF = DuploCloud Agent Framework",
        "",
        "A Python framework for building LLM-powered AI agents with:",
        "",
        "• Tool calling - Execute functions on behalf of users",
        "• Human-in-the-loop - Pause for approval on sensitive operations",
        "• Session state - Persist data across conversation turns",
        "• Streaming - Real-time token-by-token output",
        "• REST API - One-line HTTP server deployment",
        "",
        "Philosophy: Hide complexity, expose simplicity."
    ]
)

add_content_slide(
    "The Complete Picture",
    [],
    code="""from dcaf.core import Agent, serve
from dcaf.tools import tool

@tool(description="List Kubernetes pods")
def list_pods(namespace: str = "default") -> str:
    return kubectl(f"get pods -n {namespace}")

@tool(requires_approval=True, description="Delete a pod")
def delete_pod(name: str, namespace: str = "default") -> str:
    return kubectl(f"delete pod {name} -n {namespace}")

agent = Agent(tools=[list_pods, delete_pod])
serve(agent)  # Running at http://localhost:8000"""
)


# ============================================================================
# SECTION 2: THE AGENT
# ============================================================================

add_section_slide("The Agent", "Your LLM-powered assistant")

add_content_slide(
    "What is an Agent?",
    [
        "An Agent is an LLM-powered assistant that can:",
        "",
        "• Understand natural language requests",
        "• Decide which tools to use",
        "• Execute tools (with approval when needed)",
        "• Maintain conversation context",
        "• Return helpful responses",
        "",
        "The Agent is the central orchestrator of your AI workflow.",
    ]
)

add_content_slide(
    "Creating an Agent",
    ["Minimal configuration, maximum capability:"],
    code="""from dcaf.core import Agent

# Create an agent with tools
agent = Agent(
    tools=[list_pods, delete_pod, scale_deployment],
    system="You are a Kubernetes assistant for the production cluster.",
)

# Run programmatically
response = agent.run([
    {"role": "user", "content": "What pods are running?"}
])

print(response.text)
# "Here are the pods currently running: nginx, redis, api..."

# Check if approval is needed
if response.needs_approval:
    print("Approval required for:", response.pending_tools)"""
)

add_content_slide(
    "Agent Configuration",
    [
        "Constructor parameters:",
        "",
        "• tools - List of @tool decorated functions",
        "• system - System prompt for the LLM",
        "• model - LLM model ID (default: Claude 3 Sonnet)",
        "• provider - LLM provider (bedrock, anthropic, openai)",
        "• on_event - Event handlers for logging/notifications",
        "• request_interceptors - Modify requests before LLM",
        "• response_interceptors - Modify responses after LLM",
    ]
)

add_content_slide(
    "Custom Agent Functions",
    ["For complex workflows, use a function instead of Agent class:"],
    code="""from dcaf.core import Agent, AgentResult, serve

def my_agent(messages: list, context: dict) -> AgentResult:
    tenant = context.get("tenant_name")
    
    # Step 1: Classify intent
    classifier = Agent(system="Classify as: query or action")
    intent = classifier.run(messages)
    
    # Step 2: Route to appropriate handler
    if "action" in intent.text.lower():
        executor = Agent(tools=[list_pods, delete_pod])
        result = executor.run(messages)
        return AgentResult(text=result.text, pending_tools=result.pending_tools)
    
    return AgentResult(text=intent.text)

serve(my_agent)  # Works the same as Agent class"""
)


# ============================================================================
# SECTION 3: DEFINING TOOLS
# ============================================================================

add_section_slide("Defining Tools", "Functions your agent can call")

add_content_slide(
    "What is a Tool?",
    [
        "A Tool is a function that your Agent can call to:",
        "",
        "• Query external systems (databases, APIs, Kubernetes)",
        "• Perform calculations",
        "• Execute operations (create, update, delete)",
        "• Access files or resources",
        "",
        "The Agent decides which tools to use based on the user's request.",
        "",
        "You define tools with the @tool decorator."
    ]
)

add_content_slide(
    "Three Ways to Define Tool Schemas",
    [
        "Option 1: Auto-Generate (Simplest)",
        "   Schema inferred from function signature - zero config",
        "",
        "Option 2: Dict Schema (Full Control)",
        "   Pass explicit JSON Schema for validation, enums, constraints",
        "",
        "Option 3: Pydantic Model (Type-Safe)",
        "   Pass a BaseModel class for IDE support and reusability",
        "",
        "The framework normalizes all three to JSON Schema internally."
    ]
)

add_content_slide(
    "Option 1: Auto-Generated Schema",
    ["Just add type hints - DCAF does the rest:"],
    code="""from dcaf.tools import tool

@tool(description="Scale a Kubernetes deployment")
def scale_deployment(
    name: str,                    # Required string
    replicas: int,                # Required integer
    namespace: str = "default"    # Optional with default
) -> str:
    return kubectl(f"scale deployment {name} --replicas={replicas} -n {namespace}")

# Auto-generates JSON Schema:
# {
#   "type": "object",
#   "properties": {
#     "name": {"type": "string"},
#     "replicas": {"type": "integer"},
#     "namespace": {"type": "string", "default": "default"}
#   },
#   "required": ["name", "replicas"]
# }"""
)

add_content_slide(
    "Option 2: Dict Schema",
    ["Full JSON Schema control for enums, constraints, patterns:"],
    code="""@tool(
    description="Scale a deployment",
    schema={
        "type": "object",
        "properties": {
            "name": {"type": "string", "minLength": 1, "maxLength": 63},
            "replicas": {"type": "integer", "minimum": 0, "maximum": 100},
            "namespace": {
                "type": "string",
                "enum": ["default", "staging", "production"]
            }
        },
        "required": ["name", "replicas"]
    }
)
def scale_deployment(name: str, replicas: int, namespace: str = "default") -> str:
    return kubectl(f"scale deployment {name} --replicas={replicas} -n {namespace}")"""
)

add_content_slide(
    "Option 3: Pydantic Model",
    ["Type-safe with IDE autocomplete and reusability:"],
    code="""from pydantic import BaseModel, Field
from typing import Literal

class ScaleDeploymentInput(BaseModel):
    name: str = Field(..., description="Deployment name", min_length=1)
    replicas: int = Field(..., ge=0, le=100, description="Target replicas")
    namespace: Literal["default", "staging", "production"] = "default"

@tool(description="Scale a deployment", schema=ScaleDeploymentInput)
def scale_deployment(name: str, replicas: int, namespace: str = "default") -> str:
    return kubectl(f"scale deployment {name} --replicas={replicas} -n {namespace}")

# Benefits: IDE autocomplete, type checking, reusable across tools"""
)


# ============================================================================
# SECTION 4: SERVING THE AGENT
# ============================================================================

add_section_slide("Serving the Agent", "REST API in one line")

add_content_slide(
    "The serve() Function",
    ["One-line REST API server:"],
    code="""from dcaf.core import Agent, serve

agent = Agent(tools=[...])

# Development
serve(agent)  # http://0.0.0.0:8000

# Production
serve(
    agent,
    port=8000,
    host="0.0.0.0",
    workers=4,              # Multiple worker processes
    timeout_keep_alive=30,  # Match load balancer timeout
    log_level="warning",
)

# Workers: Use (2 × cpu_cores) + 1 for production
# Keep-Alive: Set to match your load balancer (AWS ALB default is 60s)"""
)

add_table_slide(
    "REST API Endpoints",
    ["Endpoint", "Method", "Description"],
    [
        ["/health", "GET", "Health check (non-blocking, immediate response)"],
        ["/api/chat", "POST", "Synchronous chat - wait for full response"],
        ["/api/chat-stream", "POST", "Streaming chat - NDJSON token stream"],
    ]
)

add_content_slide(
    "Request/Response Format",
    ["HelpDesk protocol for DuploCloud integration:"],
    code="""# Request
{
  "messages": [
    {"role": "user", "content": "List the pods in production"}
  ],
  "data": {
    "platform_context": {"tenant_name": "prod", "user_id": "alice"},
    "session": {"previous_query": "..."}
  }
}

# Response
{
  "role": "assistant",
  "content": "Here are the pods in production: nginx, redis, api...",
  "data": {
    "tool_calls": [],           # Pending approvals
    "executed_tool_calls": [],  # Completed tools
    "session": {}               # Updated session
  }
}"""
)

add_content_slide(
    "Streaming Responses",
    ["Real-time token-by-token output via /api/chat-stream:"],
    code="""# Request
curl -X POST http://localhost:8000/api/chat-stream \\
  -H "Content-Type: application/json" \\
  -d '{"messages": [{"role": "user", "content": "Explain Kubernetes"}]}'

# Response stream (NDJSON - one JSON object per line)
{"type": "text_delta", "text": "Kubernetes"}
{"type": "text_delta", "text": " is a container"}
{"type": "text_delta", "text": " orchestration platform..."}
{"type": "tool_calls", "tool_calls": [...]}
{"type": "executed_tool_calls", "executed_tool_calls": [...]}
{"type": "done"}"""
)


# ============================================================================
# SECTION 5: HUMAN-IN-THE-LOOP APPROVAL
# ============================================================================

add_section_slide("Human-in-the-Loop Approval", "Safe execution of dangerous operations")

add_content_slide(
    "Why Approval Matters",
    [
        "Dangerous operations need human oversight:",
        "",
        "• Deleting resources (pods, deployments, databases)",
        "• Modifying production configurations",
        "• Scaling infrastructure up/down",
        "• Executing shell commands",
        "• Any irreversible action",
        "",
        "Benefits:",
        "• Prevents accidental damage",
        "• Creates audit trail",
        "• Builds trust in AI-assisted workflows",
    ]
)

add_content_slide(
    "Marking Tools for Approval",
    ["Add requires_approval=True:"],
    code="""# Safe operation - executes immediately
@tool(description="List pods in a namespace")
def list_pods(namespace: str = "default") -> str:
    return kubectl(f"get pods -n {namespace}")

# Dangerous operation - pauses for approval
@tool(requires_approval=True, description="Delete a pod")
def delete_pod(name: str, namespace: str = "default") -> str:
    return kubectl(f"delete pod {name} -n {namespace}")

# Also dangerous
@tool(requires_approval=True, description="Scale deployment")
def scale_deployment(name: str, replicas: int) -> str:
    return kubectl(f"scale deployment {name} --replicas={replicas}")"""
)

add_content_slide(
    "The Approval Flow",
    [
        "1. User: \"Delete the nginx pod in production\"",
        "",
        "2. Agent identifies delete_pod requires approval",
        "",
        "3. Agent PAUSES and returns pending tool call:",
        "   {tool_calls: [{name: 'delete_pod', execute: false}]}",
        "",
        "4. User reviews and decides: Approve or Reject",
        "",
        "5. User sends approval: {execute: true}",
        "",
        "6. Agent executes tool and returns result",
    ]
)

add_content_slide(
    "Response with Pending Approval",
    ["When approval is needed, the agent returns:"],
    code="""{
  "role": "assistant",
  "content": "I'll delete the nginx pod. This requires your approval.",
  "data": {
    "tool_calls": [{
      "id": "tc_abc123",
      "name": "delete_pod",
      "input": {"name": "nginx", "namespace": "production"},
      "execute": false,
      "tool_description": "Delete a pod"
    }]
  }
}"""
)

add_content_slide(
    "Sending Approval or Rejection",
    ["Approve with execute: true, or reject with rejection_reason:"],
    code="""# Approve
{
  "messages": [...],
  "data": {
    "tool_calls": [{
      "id": "tc_abc123",
      "name": "delete_pod", 
      "input": {"name": "nginx", "namespace": "production"},
      "execute": true  # APPROVED
    }]
  }
}

# Reject
{
  "data": {
    "tool_calls": [{
      "id": "tc_abc123",
      "rejection_reason": "Wrong pod - I meant nginx-v1"  # REJECTED
    }]
  }
}"""
)

add_content_slide(
    "Programmatic Approval",
    ["Handle approvals in code:"],
    code="""response = agent.run(messages)

if response.needs_approval:
    for pending in response.pending_tools:
        print(f"Tool: {pending.name}")
        print(f"Input: {pending.input}")
        
        if confirm(f"Approve {pending.name}?"):
            pending.approve()
        else:
            pending.reject("User declined")
    
    # Continue after handling approvals
    response = agent.resume(response.conversation_id)

print(response.text)"""
)


# ============================================================================
# SECTION 6: SESSION MANAGEMENT
# ============================================================================

add_section_slide("Session Management", "Persist state across conversations")

add_content_slide(
    "What is a Session?",
    [
        "Sessions persist state across conversation turns",
        "",
        "Features:",
        "• Key-value storage",
        "• Travels with request/response automatically",
        "• Typed storage - Pydantic models & dataclasses",
        "• Auto-serialization/deserialization",
        "",
        "Use cases:",
        "• Remember user preferences",
        "• Track multi-step workflows",
        "• Store typed domain objects",
    ]
)

add_content_slide(
    "Using Session in Tools",
    ["Add session: Session parameter:"],
    code="""from dcaf.core import Session
from dcaf.tools import tool

@tool(description="Greet the user")
def greet(name: str, session: Session) -> str:
    '''Remember if we've met before.'''
    if session.get("greeted"):
        return f"Welcome back, {name}!"
    
    session.set("greeted", True)
    session.set("user_name", name)
    return f"Hello {name}, nice to meet you!"

# First call: "Hello Alice, nice to meet you!"
# Second call: "Welcome back, Alice!" """
)

add_content_slide(
    "Session API",
    [
        "Reading:",
        "• session.get(key, default) - Get value or default",
        "• session.get(key, as_type=Model) - Get as typed model",
        "• session.has(key) - Check if key exists",
        "",
        "Writing:",
        "• session.set(key, value) - Store (auto-serializes models)",
        "• session.update({...}) - Bulk update",
        "• session.delete(key) - Remove a key",
        "• session.clear() - Remove all data",
    ]
)

add_content_slide(
    "Typed Session Storage",
    ["Store Pydantic models with auto-serialization:"],
    code="""from pydantic import BaseModel
from dcaf.core import Session

class ShoppingCart(BaseModel):
    items: list[dict] = []
    discount: str | None = None

# Store - auto-serializes via model_dump()
cart = ShoppingCart(items=[{"name": "Widget", "qty": 2}])
session.set("cart", cart)

# Retrieve as typed - auto-deserializes
cart = session.get("cart", as_type=ShoppingCart)
print(cart.items[0]["name"])  # "Widget"

# Without type - returns raw dict
raw = session.get("cart")  # {"items": [...], "discount": None}"""
)

add_content_slide(
    "Session in Protocol",
    ["Session travels in data.session:"],
    code="""# Request - pass previous session
{
  "messages": [{"role": "user", "content": "Hi again"}],
  "data": {
    "session": {"greeted": true, "user_name": "Alice"}
  }
}

# Response - includes updated session
{
  "content": "Welcome back, Alice!",
  "data": {
    "session": {"greeted": true, "user_name": "Alice", "visits": 2}
  }
}

# Client stores session and sends it with next request"""
)

add_content_slide(
    "Multi-Step Workflow Example",
    [],
    code="""@tool(description="Start deployment wizard")
def start_deploy(session: Session) -> str:
    session.set("wizard_step", 1)
    return "Step 1: What service do you want to deploy?"

@tool(description="Set deployment service")  
def set_service(name: str, session: Session) -> str:
    session.set("service_name", name)
    session.set("wizard_step", 2)
    return f"Service: {name}. Step 2: How many replicas?"

@tool(description="Set replica count")
def set_replicas(count: int, session: Session) -> str:
    session.set("replicas", count)
    service = session.get("service_name")
    return f"Deploy {service} with {count} replicas? (requires approval)"

@tool(requires_approval=True, description="Execute deployment")
def execute_deploy(session: Session) -> str:
    service, replicas = session.get("service_name"), session.get("replicas")
    session.clear()  # Reset wizard
    return f"Deployed {service} with {replicas} replicas!" """
)


# ============================================================================
# SECTION 7: PROMPT CACHING (NEW)
# ============================================================================

add_section_slide("Prompt Caching", "Reduce costs by up to 90%")

add_content_slide(
    "What is Prompt Caching?",
    [
        "AWS Bedrock feature that caches static prompt content:",
        "",
        "Benefits:",
        "• Up to 90% cost reduction on cached tokens",
        "• Up to 85% latency reduction",
        "• 5-minute cache TTL (resets on each cache hit)",
        "",
        "Key Insight:",
        "Separate static instructions (cached) from dynamic context (fresh)",
        "",
        "Status: Experimental v1 - Temporary until Agno adds native support"
    ]
)

add_content_slide(
    "How It Works",
    ["DCAF places a cache checkpoint between static and dynamic parts:"],
    code="""[Static Instructions]  ← CACHED (same for all requests)
         ↓
[Cache Checkpoint]
         ↓
[Dynamic Context]      ← NOT cached (fresh each time)

Example:
• Static: "You are a K8s expert. [detailed guidelines]" ← CACHED
• Dynamic: "Tenant: acme-corp, User: alice" ← FRESH""")

add_content_slide(
    "Basic Usage",
    ["Enable caching with one parameter:"],
    code="""from dcaf.core import Agent

agent = Agent(
    system_prompt='''You are a Kubernetes expert.
    
    [Add detailed guidelines here - aim for 1024+ tokens]
    
    Guidelines:
    - Always verify namespace before operations
    - Explain what each command does
    - Ask for confirmation on destructive operations
    [... more detailed instructions ...]
    ''',
    
    tools=[list_pods, delete_pod],
    
    model_config={
        "cache_system_prompt": True  # Enable caching
    }
)""")

add_content_slide(
    "Static + Dynamic Pattern",
    ["The recommended pattern for maximum savings:"],
    code="""agent = Agent(
    # Static part - cached (same for all requests)
    system_prompt='''You are a Kubernetes expert for a multi-tenant platform.
    
    [Detailed guidelines, examples, best practices...]
    - Verify tenant context before operations
    - Follow security best practices
    - Provide clear explanations
    ''',
    
    # Dynamic part - NOT cached (changes per request)
    system_context=lambda ctx: f'''
    === CURRENT CONTEXT ===
    Tenant: {ctx.get('tenant_name')}
    Namespace: {ctx.get('k8s_namespace')}
    User: {ctx.get('user_email')}
    ''',
    
    model_config={"cache_system_prompt": True}
)""")

add_content_slide(
    "Requirements & Best Practices",
    [
        "Minimum Token Count:",
        "• Claude 3.7 Sonnet: 1024 tokens",
        "• Claude 3.5 Haiku: 2048 tokens",
        "• ~4 characters = 1 token (aim for 4000+ chars)",
        "",
        "Best Practices:",
        "1. Make static content detailed and comprehensive",
        "2. Put all variable data in system_context",
        "3. Monitor logs for cache HIT/MISS indicators",
        "4. Ensure high request volume (>1 per 5 minutes)",
    ]
)

add_content_slide(
    "Cost Savings Example",
    [
        "Scenario: 100 requests with 1500-token static prompt + 100-token dynamic",
        "",
        "Without caching:",
        "• 100 × 1600 tokens = 160,000 tokens",
        "• Cost: ~$0.48",
        "",
        "With caching:",
        "• First request: 1600 tokens (MISS)",
        "• Next 99 requests: 100 tokens each (HIT)",
        "• Total: 11,500 tokens",
        "• Cost: ~$0.035",
        "",
        "Savings: ~93% 💰"
    ]
)

add_content_slide(
    "Monitoring Cache Performance",
    ["DCAF logs cache metrics automatically:"],
    code="""# Console logs show cache status:

INFO: ✅ Cache HIT: 950 tokens reused (~90% cost reduction)
INFO: 📝 Cache MISS: 950 tokens cached for next request

# First request is always a MISS (creates cache)
Request 1: MISS (creates cache)  → Full cost
Request 2: HIT  (uses cache)     → 10% cost
Request 3: HIT  (uses cache)     → 10% cost
...
Request N: MISS (cache expired after 5 min) → Full cost

# Warnings if misconfigured:
WARNING: System prompt (~500 tokens) below minimum threshold""")


# ============================================================================
# SECTION 8: PLATFORM CONTEXT
# ============================================================================

add_section_slide("Platform Context", "Runtime environment data")

add_content_slide(
    "Platform Context",
    [
        "Automatic context extraction from requests:",
        "",
        "Available fields:",
        "• tenant_name - DuploCloud tenant",
        "• k8s_namespace - Kubernetes namespace",
        "• user_id - Requesting user",
        "• duplo_token - Authentication token",
        "• kubeconfig - Base64-encoded kubeconfig",
        "• aws_credentials - AWS credential info",
        "",
        "Access in tools via platform_context parameter.",
    ]
)

add_content_slide(
    "Using Platform Context",
    ["Add platform_context: dict parameter to your tool:"],
    code="""@tool(description="List resources for current tenant")
def list_tenant_resources(resource_type: str, platform_context: dict) -> str:
    tenant = platform_context.get("tenant_name")
    namespace = platform_context.get("k8s_namespace")
    token = platform_context.get("duplo_token")
    
    # Use context for tenant-scoped operations
    return api_call(f"/tenants/{tenant}/resources/{resource_type}")

# In custom agent functions, context is passed directly:
def my_agent(messages: list, context: dict) -> AgentResult:
    tenant = context.get("tenant_name")
    # ..."""
)


# ============================================================================
# SECTION 9: SUMMARY
# ============================================================================

add_section_slide("Summary", "Getting started with DCAF Core")

add_table_slide(
    "Key Concepts",
    ["Concept", "Description"],
    [
        ["Agent", "LLM-powered assistant - the orchestrator"],
        ["Tool", "Function the agent can call (auto/dict/Pydantic schema)"],
        ["Approval", "Human authorization for dangerous operations"],
        ["Session", "Persist state across conversation turns"],
        ["serve()", "One-line REST API server"],
        ["Platform Context", "Runtime environment (tenant, namespace)"],
    ]
)

add_content_slide(
    "Getting Started",
    [
        "Installation:",
        "   pip install git+https://github.com/duplocloud/dcaf.git",
        "",
        "Documentation:",
        "• docs/core/index.md - Core Overview",
        "• docs/core/server.md - Server Guide",
        "• docs/guides/building-tools.md - Building Tools",
        "• docs/guides/custom-agents.md - Custom Agents",
        "",
        "Examples:",
        "• examples/core_server.py - Complete example",
    ]
)

add_content_slide(
    "Key Takeaways",
    [
        "1. Agent First",
        "   The Agent is your LLM-powered orchestrator",
        "",
        "2. Tools are Functions",
        "   Define with @tool - auto, dict, or Pydantic schema",
        "",
        "3. Human-in-the-Loop",
        "   Built-in approval for dangerous operations",
        "",
        "4. Session State",
        "   Persist data across conversation turns",
        "",
        "5. One-Line Server",
        "   serve(agent) - production-ready REST API",
    ]
)

add_title_slide("Questions?", "DCAF Core - DuploCloud Agent Framework")

# Save
output_path = "docs/presentations/dcaf-core-presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
